"""
Learnify AI — Script to generate sample test fixture files.

Run once (from the /backend directory) to create the binary test assets
that test_parsers.py needs:

    python tests/create_fixtures.py

Requires: pdfplumber (already in requirements), reportlab (for PDF generation),
and python-pptx (already in requirements).
"""

from pathlib import Path

FIXTURES_DIR = Path(__file__).parent / "fixtures"
FIXTURES_DIR.mkdir(exist_ok=True)


def create_sample_pdf() -> None:
    """Create a minimal two-page PDF using reportlab."""
    pdf_path = FIXTURES_DIR / "sample.pdf"
    if pdf_path.exists():
        return

    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas

        c = canvas.Canvas(str(pdf_path), pagesize=letter)
        c.drawString(100, 750, "Learnify AI Sample PDF — Page 1")
        c.drawString(100, 720, "This is the first paragraph of the test document.")
        c.showPage()
        c.drawString(100, 750, "Learnify AI Sample PDF — Page 2")
        c.drawString(100, 720, "Machine learning is the study of algorithms.")
        c.save()
        print(f"Created {pdf_path}")
    except ImportError:
        # If reportlab is not installed, write a minimal hand-crafted PDF
        _write_minimal_pdf(pdf_path)


def _write_minimal_pdf(path: Path) -> None:
    """Write a tiny but valid PDF without any extra dependencies."""
    minimal_pdf = b"""%PDF-1.4
1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj
2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj
3 0 obj<</Type/Page/MediaBox[0 0 612 792]/Parent 2 0 R/Contents 4 0 R/Resources<</Font<</F1 5 0 R>>>>>>endobj
4 0 obj<</Length 44>>
stream
BT /F1 12 Tf 100 750 Td (Learnify AI sample text.) Tj ET
endstream
endobj
5 0 obj<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>endobj
xref
0 6
0000000000 65535 f 
0000000009 00000 n 
0000000058 00000 n 
0000000115 00000 n 
0000000266 00000 n 
0000000360 00000 n 
trailer<</Size 6/Root 1 0 R>>
startxref
441
%%EOF"""
    path.write_bytes(minimal_pdf)
    print(f"Created minimal {path}")


def create_sample_pptx() -> None:
    """Create a minimal two-slide PPTX."""
    pptx_path = FIXTURES_DIR / "sample.pptx"
    if pptx_path.exists():
        return

    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    blank_layout = prs.slide_layouts[5]

    for i, title_text in enumerate(
        ["Introduction to Learnify AI", "Key Features"], start=1
    ):
        slide = prs.slides.add_slide(blank_layout)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        tf = txBox.text_frame
        tf.text = f"Slide {i}: {title_text}"
        tf.paragraphs[0].runs[0].font.size = Pt(24)

        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(3))
        tf2 = txBox2.text_frame
        tf2.text = (
            f"This slide covers topic {i}. "
            "Adaptive learning uses AI to personalise education."
        )

    prs.save(str(pptx_path))
    print(f"Created {pptx_path}")


def create_sample_txt() -> None:
    """Create a plain text file with multiple paragraphs."""
    txt_path = FIXTURES_DIR / "sample.txt"
    if txt_path.exists():
        return

    content = (
        "Paragraph one: Learnify AI is a multimodal learning platform.\n"
        "It combines RAG, adaptive quizzes, and emotion detection.\n"
        "\n"
        "Paragraph two: The backend is built with FastAPI and async MongoDB.\n"
        "Embeddings are generated using sentence-transformers.\n"
        "\n"
        "Paragraph three: FAISS enables fast nearest-neighbour search over chunks.\n"
        "This allows the system to retrieve relevant content in milliseconds.\n"
    )
    txt_path.write_text(content, encoding="utf-8")
    print(f"Created {txt_path}")


if __name__ == "__main__":
    create_sample_pdf()
    create_sample_pptx()
    create_sample_txt()
    print("All fixtures created in", FIXTURES_DIR)
