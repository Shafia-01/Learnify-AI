"""
Learnify AI — PowerPoint Parser.

Extracts text from PPTX presentation files slide by slide using python-pptx.
Text is gathered from every text frame on each slide and joined with newlines.
"""

import logging
from pathlib import Path
from typing import Any, Dict, List

from pptx import Presentation

logger = logging.getLogger(__name__)


def parse_ppt(file_path: str) -> List[Dict[str, Any]]:
    """
    Extract text from a PowerPoint (.pptx) file, slide by slide.

    Iterates every shape on each slide that has a text frame, concatenates
    all non-empty paragraph texts, and produces one output dict per slide.

    Args:
        file_path: Absolute or relative path to the PPTX file.

    Returns:
        A list of dicts, each containing:
            - ``text``  — Concatenated text from all text frames on the slide.
            - ``slide`` — The 1-based slide number.
            - ``source_file`` — The filename of the source PPTX.
    """
    path = Path(file_path)
    source_file = path.name
    results: List[Dict[str, Any]] = []

    prs = Presentation(str(path))

    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts: List[str] = []

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        slide_texts.append(line)

        full_text = "\n".join(slide_texts)

        if full_text.strip():
            results.append(
                {
                    "text": full_text.strip(),
                    "slide": slide_num,
                    "source_file": source_file,
                }
            )
        else:
            logger.warning(
                "Slide %d of '%s' has no extractable text — skipping.",
                slide_num,
                source_file,
            )

    return results
